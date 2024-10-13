from pptx import Presentation
from FramePractice import Ui_Form#(注意這邊名稱要看原來建立class是甚麼類別，例如有Dialogue,or Form)
from PyQt5.QtWidgets import *
from PyQt5.QtCore import *
from PyQt5.QtGui import *
import sys

app = QApplication(sys.argv)#創建應用程式的主對象，負責控制應用的啟動和結束。
#Pyuic5電腦會去找python.exe接著把-x,檔名,-o,檔名丟給Pyuic5等四個參數
widget = QWidget()#創建一個空白的主窗口（視窗）
ui = Ui_Form()
ui.setupUi(widget)
def extract_notes_from_pptx(pptx_path, output_txt_path):
    notes_content = []#負責儲存所有備忘錄文字，準備被return
    # 打開PPTX檔案
    print(f"正在打開PPTX檔案: {pptx_path}")
    presentation = Presentation(pptx_path)
    
    # 開啟輸出文字檔案
    print(f"正在打開文字檔案進行寫入: {output_txt_path}")
    with open(output_txt_path, 'w', encoding='utf-8') as file:
        # 遍歷每張投影片
        for i, slide in enumerate(presentation.slides):
            print(f"正在處理投影片 {i+1}...")
            # 提取投影片備忘錄（notes_slide）
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text.strip()
                # 確認備忘錄是否有內容
                if notes_text:
                    notes_content.append(f"Slide {i + 1} Notes:\n{notes_text}\n")
                    file.write(f"Slide {i+1} Notes:\n")
                    file.write(notes_text + '\n\n')
                    #print(f"投影片 {i+1} 備忘錄: {notes_text}")
                else:
                    notes_content.append(f"Slide {i + 1} has empty notes.\n")
                    file.write(f"Slide {i+1} has empty notes.\n\n")
                    #print(f"投影片 {i+1} 備忘錄為空")
            else:
                notes_content.append(f"Slide {i + 1} has no notes.\n")
                # 若無備忘錄，提示無備忘錄訊息
                file.write(f"Slide {i+1} has no notes.\n\n")
                #print(f"投影片 {i+1} 沒有備忘錄")
    return ''.join(notes_content)  # 要有return才能夠導出來給textBrowser顯示
    print(f"備忘錄已提取並儲存至 {output_txt_path}")

# 正確處理路徑問題，使用正斜杠
pptx_path = "D:/工作區/QT_Project/PPTX_Project/test.pptx"
#修改成任意檔名

output_txt_path = "D:/工作區/QT_Project/PPTX_Project/test.txt"

# 呼叫函數
#extract_notes_from_pptx(pptx_path, output_txt_path)

def browse_file():
        # 打開檔案選擇對話框，過濾 pptx 檔案
        pptx_path, _ = QFileDialog.getOpenFileName(widget,"選擇PPTX檔案", "", "PPTX Files (*.pptx)")
        
        if pptx_path:
            # 更新檔案路徑文字框
            ui.label.setText(pptx_path)
            
            # 設置輸出的文字檔案路徑（這裡可以動態生成或指定）
            output_txt_path = "output.txt"
            # 提取備忘錄並顯示在 TextBrowser 中
            notes = extract_notes_from_pptx(pptx_path, output_txt_path)
            ui.textBrowser.setText(notes)
            
def resetContext():
    ui.textBrowser.setText("")
    
ui.pushButton.clicked.connect(browse_file)
ui.pushButton_2.clicked.connect(resetContext)
widget.show()
app.exec_()


